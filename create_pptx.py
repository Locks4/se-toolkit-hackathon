"""
Generate PowerPoint Presentation for Goal Tracker
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
PRIMARY = RGBColor(0x66, 0x7E, 0xEA)
SECONDARY = RGBColor(0x76, 0x4B, 0xA2)
ACCENT = RGBColor(0x4C, 0xAF, 0x50)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF5, 0xF7, 0xFA)
ORANGE = RGBColor(0xF5, 0x7C, 0x00)

def add_gradient_background(slide, color1, color2=None):
    """Add solid background (PowerPoint doesn't support gradients in code easily)"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color1

def add_shape(slide, left, top, width, height, fill_color, line_color=None):
    """Add a rectangle shape"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    """Add a rounded rectangle"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT):
    """Add a text box"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=DARK_GRAY):
    """Add a bulleted list"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return txBox

def add_icon_text(slide, left, top, icon, text, font_size=16, color=DARK_GRAY):
    """Add icon with text"""
    add_text_box(slide, left, top, Inches(0.5), Inches(0.5), icon, font_size=font_size+4, color=color, bold=True)
    add_text_box(slide, left + Inches(0.5), top, Inches(6), Inches(0.5), text, font_size=font_size, color=color)

# ========================
# SLIDE 1: TITLE SLIDE
# ========================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
add_gradient_background(slide1, PRIMARY)

# Add decorative shapes
add_shape(slide1, Inches(0), Inches(0), prs.slide_width, Inches(7.5), PRIMARY)

# Center icon
add_text_box(slide1, Inches(5.6), Inches(1.5), Inches(2), Inches(1.5), "🎯", 
             font_size=72, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Title
add_text_box(slide1, Inches(1), Inches(3), Inches(11.333), Inches(1), "Goal Tracker",
             font_size=54, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide1, Inches(2), Inches(4), Inches(9.333), Inches(0.8),
             "A Duolingo-style goal tracking app with streaks, gems, and calendar visualization",
             font_size=24, color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)

# Author info box
author_box = add_rounded_rect(slide1, Inches(4), Inches(5.2), Inches(5.333), Inches(1.8), 
                               RGBColor(0x55, 0x66, 0xDD))
author_box.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
author_box.line.width = Pt(2)

add_text_box(slide1, Inches(4.5), Inches(5.3), Inches(4.333), Inches(0.4), "Alexey",
             font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide1, Inches(4.5), Inches(5.7), Inches(4.333), Inches(0.35), "📧 a.petrov@innopolis.university",
             font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide1, Inches(4.5), Inches(6.05), Inches(4.333), Inches(0.35), "👥 Group: DSAI-05",
             font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide1, Inches(4.5), Inches(6.4), Inches(4.333), Inches(0.35), "🐙 github.com/Locks4",
             font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)

# ========================
# SLIDE 2: CONTEXT
# ========================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide2, WHITE)

# Title
add_text_box(slide2, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8), "🎯 Context",
             font_size=40, color=PRIMARY, bold=True)

# End Users Box
box1 = add_rounded_rect(slide2, Inches(0.5), Inches(1.3), Inches(12.333), Inches(1.4), 
                        RGBColor(0xF0, 0xF4, 0xFF), PRIMARY)
add_icon_text(slide2, Inches(0.8), Inches(1.4), "👤", "End Users", 20, PRIMARY)
add_text_box(slide2, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.8),
             "Students and professionals who want to build consistent daily habits and track personal goals with gamified motivation.",
             font_size=18, color=DARK_GRAY)

# Problem Box
box2 = add_rounded_rect(slide2, Inches(0.5), Inches(3), Inches(12.333), Inches(1.4),
                        RGBColor(0xF0, 0xF4, 0xFF), PRIMARY)
add_icon_text(slide2, Inches(0.8), Inches(3.1), "❌", "Problem", 20, PRIMARY)
add_text_box(slide2, Inches(0.8), Inches(3.5), Inches(11.5), Inches(0.8),
             "People struggle to maintain consistency with their goals. Without proper tracking and motivation, it's easy to lose track of progress and break streaks.",
             font_size=18, color=DARK_GRAY)

# Solution Box
box3 = add_rounded_rect(slide2, Inches(0.5), Inches(4.7), Inches(12.333), Inches(2.2),
                        RGBColor(0xE8, 0xF5, 0xE9), ACCENT)
add_icon_text(slide2, Inches(0.8), Inches(4.8), "💡", "Solution", 20, ACCENT)
add_text_box(slide2, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.5),
             "Goal Tracker combines Duolingo's addictive streak system with personal goal tracking to help users build consistent habits through gamification, visual progress tracking, and accountability.",
             font_size=20, color=DARK_GRAY, bold=True)

# ========================
# SLIDE 3: IMPLEMENTATION
# ========================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide3, WHITE)

# Title
add_text_box(slide3, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8), "⚙️ Implementation",
             font_size=40, color=PRIMARY, bold=True)

# Tech Stack Column
add_rounded_rect(slide3, Inches(0.5), Inches(1.2), Inches(5.8), Inches(2.5), LIGHT_GRAY)
add_text_box(slide3, Inches(0.8), Inches(1.3), Inches(5), Inches(0.5), "🔧 Tech Stack",
             font_size=24, color=PRIMARY, bold=True)
add_bullet_list(slide3, Inches(0.8), Inches(1.8), Inches(5), Inches(2),
               ["Backend: Python FastAPI", "Frontend: React + TypeScript", 
                "Database: SQLite (SQLAlchemy)", "Auth: JWT tokens", "Deployment: Docker Compose"],
               font_size=16)

# Architecture Column
add_rounded_rect(slide3, Inches(6.8), Inches(1.2), Inches(6), Inches(2.5), LIGHT_GRAY)
add_text_box(slide3, Inches(7.1), Inches(1.3), Inches(5.5), Inches(0.5), "📦 Architecture",
             font_size=24, color=PRIMARY, bold=True)
add_bullet_list(slide3, Inches(7.1), Inches(1.8), Inches(5.5), Inches(2),
               ["RESTful API backend", "Single-page React frontend",
                "Date-specific goal tracking", "Per-day subgoal completion", "Calendar heatmap visualization"],
               font_size=16)

# Versions Title
add_text_box(slide3, Inches(0.5), Inches(4), Inches(12), Inches(0.6), "🚀 Versions",
             font_size=28, color=SECONDARY, bold=True)

# V1 Column
add_rounded_rect(slide3, Inches(0.5), Inches(4.6), Inches(5.8), Inches(2.6), RGBColor(0xE3, 0xF2, 0xFD))
add_text_box(slide3, Inches(0.8), Inches(4.7), Inches(5), Inches(0.5), "Version 1",
             font_size=22, color=RGBColor(0x19, 0x76, 0xD2), bold=True)
add_bullet_list(slide3, Inches(0.8), Inches(5.2), Inches(5), Inches(2),
               ["Basic goal CRUD", "Streak tracking", "Gems system (5 gems per freeze)",
                "Calendar view", "User authentication"],
               font_size=15)

# V2 Column
add_rounded_rect(slide3, Inches(6.8), Inches(4.6), Inches(6), Inches(2.6), RGBColor(0xE8, 0xF5, 0xE9))
add_text_box(slide3, Inches(7.1), Inches(4.7), Inches(5.5), Inches(0.5), "Version 2",
             font_size=22, color=RGBColor(0x38, 0x8E, 0x3C), bold=True)
add_bullet_list(slide3, Inches(7.1), Inches(5.2), Inches(5.5), Inches(2),
               ["Subgoals with daily tracking", "One-time scheduled goals", "Day navigation",
                "Beautiful UI modals", "Docker deployment", "Desktop app mode"],
               font_size=15)

# TA Feedback Box
add_rounded_rect(slide3, Inches(0.5), Inches(4.6), Inches(0), Inches(0), RGBColor(0xFF, 0xF3, 0xCD), ORANGE)
add_text_box(slide3, Inches(0.5), Inches(4.6), Inches(12), Inches(0.5), "📝 Not Created Yet",
             font_size=18, color=ORANGE, bold=True)

# ========================
# SLIDE 4: DEMO
# ========================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide4, WHITE)

# Title
add_text_box(slide4, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8), "🎬 Demo",
             font_size=40, color=PRIMARY, bold=True)

# Video placeholder
video_box = add_rounded_rect(slide4, Inches(0.5), Inches(1.2), Inches(12.333), Inches(2.5), 
                              RGBColor(0x00, 0x00, 0x00))
add_text_box(slide4, Inches(1), Inches(1.5), Inches(11.333), Inches(0.8), "🎥",
             font_size=60, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide4, Inches(1), Inches(2.3), Inches(11.333), Inches(0.6), "Demo Video",
             font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide4, Inches(1), Inches(2.8), Inches(11.333), Inches(0.5), 
             "Version 2 demonstration with voice-over",
             font_size=18, color=RGBColor(0xCC, 0xCC, 0xCC), alignment=PP_ALIGN.CENTER)

# Demo video note
add_rounded_rect(slide4, Inches(0.5), Inches(3.9), Inches(12.333), Inches(0.8), RGBColor(0xFF, 0xF3, 0xCD), ORANGE)
add_text_box(slide4, Inches(0.8), Inches(4), Inches(11.5), Inches(0.6),
             "✅ Demo video has been added to the repository!",
             font_size=18, color=RGBColor(0x85, 0x64, 0x04), bold=True)

# Feature checklist
add_text_box(slide4, Inches(0.5), Inches(4.9), Inches(12), Inches(0.5), "Features Demonstrated:",
             font_size=24, color=PRIMARY, bold=True)

features = [
    ("✓", "User registration & login"),
    ("✓", "Adding daily & one-time goals"),
    ("✓", "Completing goals & earning gems"),
    ("✓", "Subgoals with daily tracking"),
    ("✓", "Calendar heatmap navigation"),
    ("✓", "Streak freezes & stats")
]

for i, (check, feature) in enumerate(features):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(6)
    y = Inches(5.4) + row * Inches(0.5)
    add_text_box(slide4, x, y, Inches(0.4), Inches(0.4), check, 
                 font_size=18, color=ACCENT, bold=True)
    add_text_box(slide4, x + Inches(0.5), y, Inches(5), Inches(0.4), feature,
                 font_size=16, color=DARK_GRAY)

# ========================
# SLIDE 5: LINKS
# ========================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide5, WHITE)

# Title
add_text_box(slide5, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8), "🔗 Links & Resources",
             font_size=40, color=PRIMARY, bold=True)

# GitHub Card
add_rounded_rect(slide5, Inches(0.5), Inches(1.3), Inches(5.8), Inches(3.5), LIGHT_GRAY)
add_text_box(slide5, Inches(0.8), Inches(1.5), Inches(5), Inches(0.5), "📂 GitHub Repository",
             font_size=24, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

# QR Code placeholder
qr1 = add_rounded_rect(slide5, Inches(2.3), Inches(2.2), Inches(2.5), Inches(2.5), WHITE)
add_text_box(slide5, Inches(2.5), Inches(2.8), Inches(2), Inches(0.4), "QR Code",
             font_size=16, color=DARK_GRAY, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide5, Inches(2.5), Inches(3.2), Inches(2), Inches(0.4), "Scan me!",
             font_size=14, color=RGBColor(0x99, 0x99, 0x99), alignment=PP_ALIGN.CENTER)

add_text_box(slide5, Inches(0.8), Inches(4.9), Inches(5), Inches(0.4),
             "github.com/Locks4/se-toolkit-hackathon",
             font_size=14, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

# Deployed App Card
add_rounded_rect(slide5, Inches(6.8), Inches(1.3), Inches(6), Inches(3.5), LIGHT_GRAY)
add_text_box(slide5, Inches(7.1), Inches(1.5), Inches(5.5), Inches(0.5), "🌐 Deployed Application",
             font_size=24, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

# QR Code placeholder
qr2 = add_rounded_rect(slide5, Inches(8.6), Inches(2.2), Inches(2.5), Inches(2.5), WHITE)
add_text_box(slide5, Inches(8.8), Inches(2.8), Inches(2), Inches(0.4), "QR Code",
             font_size=16, color=DARK_GRAY, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide5, Inches(8.8), Inches(3.2), Inches(2), Inches(0.4), "Scan me!",
             font_size=14, color=RGBColor(0x99, 0x99, 0x99), alignment=PP_ALIGN.CENTER)

add_text_box(slide5, Inches(7.1), Inches(4.9), Inches(5.5), Inches(0.4),
             "localhost:80 (docker compose up -d)",
             font_size=14, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

# Footer info box
add_rounded_rect(slide5, Inches(0.5), Inches(5.2), Inches(12.333), Inches(2), LIGHT_GRAY)
add_text_box(slide5, Inches(0.8), Inches(5.4), Inches(11.5), Inches(0.4),
             "📖 Documentation: README.md | DEPLOYMENT.md | API Docs at /docs",
             font_size=18, color=DARK_GRAY, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide5, Inches(0.8), Inches(5.9), Inches(11.5), Inches(0.4),
             "Quick Deploy:",
             font_size=16, color=DARK_GRAY, bold=True)
add_text_box(slide5, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.8),
             "git clone https://github.com/Locks4/se-toolkit-hackathon.git",
             font_size=16, color=PRIMARY)
add_text_box(slide5, Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.4),
             "cd se-toolkit-hackathon && docker compose up -d",
             font_size=16, color=PRIMARY)

# Save presentation
output_path = r"e:\Git-projects\se-toolkit-hackathon\se-toolkit-hackathon\Goal_Tracker_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
